import sys
import os
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

import fitz
from clean_text import clean_text
from chunk_text import chunk_text
try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None

def text_extraction(file_path):
    all_chunks = []
    
    if file_path.lower().endswith('.txt'):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                raw_text = f.read()
            if '\x00' in raw_text:
                with open(file_path, 'r', encoding='utf-16') as f:
                    raw_text = f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                raw_text = f.read()
                
        cleaned_text = clean_text(raw_text)
        chunks = chunk_text(cleaned_text)
        
        for idx, single_chunk in enumerate(chunks):
            all_chunks.append({
                "pages": 1, 
                "chunk_id": idx,
                "text": single_chunk
            })
    elif file_path.lower().endswith('.docx'):
        if docx:
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            cleaned_text = clean_text("\n".join(full_text))
            chunks = chunk_text(cleaned_text)
            for idx, single_chunk in enumerate(chunks):
                all_chunks.append({
                    "pages": 1,
                    "chunk_id": idx,
                    "text": single_chunk
                })
        else:
            print("python-docx not installed.")
    elif file_path.lower().endswith('.pptx'):
        if pptx:
            prs = pptx.Presentation(file_path)
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                cleaned_text = clean_text("\n".join(slide_text))
                chunks = chunk_text(cleaned_text)
                for idx, single_chunk in enumerate(chunks):
                    all_chunks.append({
                        "pages": i + 1,
                        "chunk_id": idx,
                        "text": single_chunk
                    })
        else:
            print("python-pptx not installed.")
    elif file_path.lower().endswith('.pdf'):
        # Assume PDF
        doc = fitz.open(file_path)
        for i in range(len(doc)):
            page = doc[i]
            raw_text = page.get_text("text")
            cleaned_text = clean_text(raw_text)
            chunks = chunk_text(cleaned_text)
    
            for idx, single_chunk in enumerate(chunks):
                all_chunks.append({
                    "pages": i + 1,
                    "chunk_id": idx,
                    "text": single_chunk
                })
    else:
        print(f"Unsupported file type: {file_path}")

    return all_chunks


if __name__ == "__main__":
    chunks = text_extraction("data/indian stock market.pdf")
    print(f"Extracted {len(chunks)} chunks")
